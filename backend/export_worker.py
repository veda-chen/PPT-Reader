"""
高保真导出 worker — 独立子进程运行。

流程：
  1. PowerPoint COM 打开 PPT，SaveAs 导出为 PDF（保真度最高）
  2. PyMuPDF 打开该 PDF：
     - 每页渲染为高清 PNG（背景图）
     - 提取每个文字 span 的精确 bbox（透明文字层用，与 PNG 同坐标系，像素级对齐）
  3. 输出 render/slide{N}.png 和 render/spans.json

PNG 与文字 bbox 来自同一个 PDF，因此选中/高亮的透明文字层能与可见 PNG 完美对齐。

用法:
  python export_worker.py <pptx_path> <out_dir> [zoom]
stdout 输出 JSON: {"ok": true, "count": N, "page_width": W_pt, "page_height": H_pt}
"""
import sys
import os
import json

PP_SAVE_AS_PDF = 32  # ppSaveAsPDF


# ── 文字提取辅助函数 ──────────────────────────────────────

def _extract_spans_dict(page) -> list[dict]:
    """
    从 PyMuPDF dict 提取文字 span（结构化 block→line→span）。
    dict 的 span 直接带 text/size/bbox（rawdict 的 span 无 text 字段，文字在 chars 里，
    早期误用 rawdict 导致该通道恒返回空——见 git 历史）。
    返回结构: [{text, x, y, w, h, size, block, line, span}, ...]
    """
    spans = []
    try:
        d = page.get_text("dict")
    except Exception:
        return spans

    for b_idx, block in enumerate(d.get("blocks", [])):
        if block.get("type") != 0:  # 0=文字块，跳过图片/路径块
            continue
        for l_idx, line in enumerate(block.get("lines", [])):
            for s_idx, sp in enumerate(line.get("spans", [])):
                text = sp.get("text", "")
                # 放宽过滤：仅跳过完全为空的 span（保留单个空格等）
                if not text:
                    continue
                x0, y0, x1, y1 = sp["bbox"]
                spans.append({
                    "text": text,
                    "x": round(x0, 2),
                    "y": round(y0, 2),
                    "w": round(x1 - x0, 2),
                    "h": round(y1 - y0, 2),
                    "size": round(sp.get("size", 12), 2),
                    "block": b_idx,
                    "line": l_idx,
                    "span": s_idx,
                })
    return spans


def _extract_spans_words(page) -> list[dict]:
    """
    从 PyMuPDF words 平铺提取单词级 span。
    words 返回 [(x0,y0,x1,y1,word,block_no,line_no,word_no), ...]，
    可捕获 rawdict 层级中因结构异常而遗漏的文字。
    字段映射: block_no→block, line_no→line, word_no→span
    """
    spans = []
    try:
        words = page.get_text("words")
    except Exception:
        return spans

    for w in words:
        x0, y0, x1, y1, text, block_no, line_no, word_no = w
        if not text or not text.strip():
            continue
        spans.append({
            "text": text,
            "x": round(x0, 2),
            "y": round(y0, 2),
            "w": round(x1 - x0, 2),
            "h": round(y1 - y0, 2),
            "size": round(y1 - y0, 2),  # words 无字号，用 bbox 高度估算
            "block": block_no,
            "line": line_no,
            "span": word_no,
        })
    return spans


def _has_overlap(span: dict, existing: list[dict], threshold: float = 0.5) -> bool:
    """
    判断 span 与 existing 列表中的任意 span 是否有显著 bbox 重叠。
    重叠面积 / min(span面积, existing面积) > threshold 则认为重复。
    """
    ax1, ay1 = span["x"], span["y"]
    ax2, ay2 = ax1 + span["w"], ay1 + span["h"]
    area_a = span["w"] * span["h"]
    if area_a <= 0:
        return False

    for b in existing:
        bx1, by1 = b["x"], b["y"]
        bx2, by2 = bx1 + b["w"], by1 + b["h"]
        # 计算交集
        ix = max(0, min(ax2, bx2) - max(ax1, bx1))
        iy = max(0, min(ay2, by2) - max(ay1, by1))
        inter = ix * iy
        if inter > 0:
            area_b = b["w"] * b["h"]
            min_area = min(area_a, area_b)
            if min_area > 0 and inter / min_area > threshold:
                return True
    return False


# 合成 span 的 block 号基址：远高于真实 block 号（个位/十位），保证不冲突且排序在后
_SYNTH_BLOCK_BASE = 100000


def _norm_text(t: str) -> str:
    """去除所有空白并小写，用于判断某段文字是否已存在于 PDF 文字层。"""
    return "".join((t or "").split()).lower()


def _recover_rasterized_spans(pptx_path: str, slides_spans: list, page_width: float, page_height: float) -> None:
    """
    用 python-pptx 把被 PowerPoint 栅格化（从 PDF 文字层消失）的文本框文字补回 slides_spans。

    带特效的文本框导出 PDF 时被转为位图，PyMuPDF 提取不到。但这些文字在 python-pptx
    里仍是真实文本，且 PPT 用 EMU 坐标（EMU/12700=pt），与 spans.json 同坐标系。
    对每页顶层文本框，凡其段落文字未出现在已提取文字层中者，按文本框 bbox 合成 span 注入。

    就地修改 slides_spans（每页追加合成 span）。v1 仅处理顶层 shape，GROUP 内暂不恢复。
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pres = Presentation(pptx_path)
    slide_w_emu = pres.slide_width or 12192000
    slide_h_emu = pres.slide_height or 6858000

    emu_to_pt = 1.0 / 12700.0
    # 防御：PDF 页尺寸与 pptx 点尺寸不一致时做归一缩放（正常 sx=sy=1）
    sx = page_width / (slide_w_emu * emu_to_pt) if slide_w_emu else 1.0
    sy = page_height / (slide_h_emu * emu_to_pt) if slide_h_emu else 1.0

    slides = list(pres.slides)
    for i, slide in enumerate(slides):
        if i >= len(slides_spans):
            break
        existing = slides_spans[i]
        extracted_norm = _norm_text("".join(sp.get("text", "") for sp in existing))

        synth_idx = 0
        group_skipped = 0
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_skipped += 1
                    continue
                if not shape.has_text_frame:
                    continue
                if shape.left is None or shape.top is None or not shape.width or not shape.height:
                    continue
            except Exception:
                continue

            paras = shape.text_frame.paragraphs
            n = max(len(paras), 1)
            sh_x = shape.left * emu_to_pt * sx
            sh_y = shape.top * emu_to_pt * sy
            sh_w = shape.width * emu_to_pt * sx
            sh_h = shape.height * emu_to_pt * sy
            row_h = sh_h / n

            for p_idx, para in enumerate(paras):
                ptext = "".join(r.text or "" for r in para.runs) or para.text or ""
                pnorm = _norm_text(ptext)
                if not pnorm:
                    continue
                if pnorm in extracted_norm:
                    continue  # 该段是矢量文字，已在文字层中，无需补

                # 字号：优先取首个 run 的真实字号，否则按行高估算
                size_pt = None
                for r in para.runs:
                    if r.font is not None and r.font.size:
                        size_pt = r.font.size * emu_to_pt
                        break
                if not size_pt:
                    size_pt = max(8.0, min(row_h * 0.6, 40.0))

                slides_spans[i].append({
                    "text": ptext,
                    "x": round(sh_x, 2),
                    "y": round(sh_y + p_idx * row_h, 2),
                    "w": round(sh_w, 2),
                    "h": round(row_h, 2),
                    "size": round(size_pt, 2),
                    "block": _SYNTH_BLOCK_BASE + synth_idx,
                    "line": 0,
                    "span": 0,
                    "synthetic": True,
                })
                synth_idx += 1

        if synth_idx:
            print(f"[恢复] 第{i+1}页: +{synth_idx} 段栅格化文本")
        if group_skipped:
            print(f"[恢复] 第{i+1}页: 跳过 {group_skipped} 个组合(GROUP)内文本(v1 未处理)")


def export(pptx_path: str, out_dir: str, zoom: float = 2.0) -> dict:
    import win32com.client
    import pythoncom
    import fitz  # PyMuPDF

    pptx_path = os.path.abspath(pptx_path)
    out_dir = os.path.abspath(out_dir)
    os.makedirs(out_dir, exist_ok=True)
    pdf_path = os.path.join(out_dir, "slides.pdf")

    # ── 1. COM 导出 PDF ──────────────────────────────
    pythoncom.CoInitialize()
    app = None
    pres = None
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(pptx_path, ReadOnly=True, WithWindow=False)
        pres.SaveAs(pdf_path, PP_SAVE_AS_PDF)
    finally:
        try:
            if pres is not None:
                pres.Close()
        except Exception:
            pass
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()

    if not os.path.exists(pdf_path):
        raise RuntimeError("PDF 导出失败")

    # ── 2. PyMuPDF 渲染 PNG + 提取文字层 ──────────────
    doc = fitz.open(pdf_path)
    page_width = doc[0].rect.width if len(doc) else 960
    page_height = doc[0].rect.height if len(doc) else 540

    slides_spans = []  # 每页一个 span 列表
    mat = fitz.Matrix(zoom, zoom)

    for i, page in enumerate(doc):
        # 渲染 PNG
        pix = page.get_pixmap(matrix=mat, alpha=False)
        pix.save(os.path.join(out_dir, f"slide{i}.png"))

        # ── 提取文字 span ──────────────────────────────
        # 策略：双通道提取
        #   通道 A — dict：结构化层级提取（block→line→span），含真实字号/字体信息
        #   通道 B — words：平铺式单词提取，可捕获 dict 可能遗漏的边缘文字
        # 合并后按 bbox 去重，结构化的 dict 条目优先
        spans_raw = _extract_spans_dict(page)
        spans_w = _extract_spans_words(page)

        # 合并去重：如果 words 中的 span 与 dict 中的 span bbox 高度重叠则跳过
        merged = list(spans_raw)
        for sw in spans_w:
            if not _has_overlap(sw, spans_raw):
                merged.append(sw)

        print(f"[提取] 第{i+1}页: dict={len(spans_raw)} words={len(spans_w)} 合并={len(merged)}")
        slides_spans.append(merged)

    doc.close()

    # ── 2.5 恢复被 PowerPoint 栅格化的文本框文字 ───────
    # 带特效（发光/阴影/渐变/图片填充等）的文本框导出 PDF 时会被转成位图，
    # 从 PDF 文字层消失。这些文字在 python-pptx 里仍是真实文本，用其 bbox 补回。
    try:
        _recover_rasterized_spans(pptx_path, slides_spans, page_width, page_height)
    except Exception as e:
        print(f"[恢复] 跳过（{type(e).__name__}: {e}）")

    # ── 3. 保存文字层 JSON ───────────────────────────
    with open(os.path.join(out_dir, "spans.json"), "w", encoding="utf-8") as f:
        json.dump({
            "page_width": page_width,
            "page_height": page_height,
            "slides": slides_spans,
        }, f, ensure_ascii=False)

    return {
        "ok": True,
        "count": len(slides_spans),
        "page_width": page_width,
        "page_height": page_height,
    }


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(json.dumps({"ok": False, "error": "用法: export_worker.py <pptx> <out_dir> [zoom]"}))
        sys.exit(1)

    pptx_path = sys.argv[1]
    out_dir = sys.argv[2]
    zoom = float(sys.argv[3]) if len(sys.argv) > 3 else 2.0

    try:
        result = export(pptx_path, out_dir, zoom)
        print(json.dumps(result))
    except Exception as e:
        print(json.dumps({"ok": False, "error": f"{type(e).__name__}: {e}"}))
        sys.exit(1)
