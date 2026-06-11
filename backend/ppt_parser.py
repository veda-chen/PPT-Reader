"""
PPT 文件解析模块 — 使用 python-pptx 提取文本、图片和元信息。
"""
import uuid
import os
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from paths import data_dir

UPLOADS_DIR = os.path.join(data_dir(), "uploads")


def save_uploaded_ppt(file_content: bytes, original_filename: str) -> tuple[str, str]:
    """
    保存上传的 .pptx 文件，返回 (ppt_id, file_path)。
    """
    ppt_id = uuid.uuid4().hex[:12]
    ppt_dir = os.path.join(UPLOADS_DIR, ppt_id)
    os.makedirs(ppt_dir, exist_ok=True)

    file_path = os.path.join(ppt_dir, "original.pptx")
    with open(file_path, "wb") as f:
        f.write(file_content)

    return ppt_id, file_path


def extract_images(pres: Presentation, ppt_id: str):
    """
    提取 PPT 中所有图片并保存到 uploads/<ppt_id>/images/ 目录。
    返回 dict: {shape_id: relative_image_path}
    """
    images_dir = os.path.join(UPLOADS_DIR, ppt_id, "images")
    os.makedirs(images_dir, exist_ok=True)

    image_map = {}
    for slide_idx, slide in enumerate(pres.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    ext = shape.image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    filename = f"slide{slide_idx}_shape{shape.shape_id}.{ext}"
                    filepath = os.path.join(images_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(image_bytes)
                    image_map[str(shape.shape_id)] = filename
                except Exception:
                    continue
    return image_map


def parse_ppt(file_path: str, ppt_id: str) -> dict:
    """
    打开 PPT 文件，提取元信息和各页文本。
    返回: { title, slide_count, slide_texts, image_map }
    """
    pres = Presentation(file_path)

    # 提取图片
    image_map = extract_images(pres, ppt_id)

    slide_texts = []
    for slide in pres.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = []
                    for cell in row.cells:
                        ct = cell.text.strip()
                        if ct:
                            row_texts.append(ct)
                    if row_texts:
                        texts.append(" | ".join(row_texts))
        slide_texts.append("\n".join(texts))

    # 尝试提取标题（第一页第一个非空文本）
    title = ""
    if slide_texts:
        lines = slide_texts[0].split("\n")
        title = lines[0][:200] if lines else ""

    return {
        "title": title,
        "slide_count": len(pres.slides),
        "slide_texts": slide_texts,
        "image_map": image_map,
    }


def get_presentation(pres: Presentation):
    """返回 python-pptx 的 Presentation 对象（用于后续渲染）。"""
    return pres


def load_presentation(ppt_id: str):
    """从磁盘加载已上传的 PPT 文件，返回 Presentation 对象。"""
    file_path = os.path.join(UPLOADS_DIR, ppt_id, "original.pptx")
    if not os.path.exists(file_path):
        return None
    return Presentation(file_path)
