"""
PPT → HTML 核心转换模块。
将 python-pptx 的 slide 对象渲染为可选中文本的 HTML 字符串。
每个 <span> 携带 data-shape-id / data-para-idx / data-run-idx 属性，
用于高亮系统的 DOM 定位。
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
import os

# ── 常量 ─────────────────────────────────────────────────
DISPLAY_WIDTH_PX = 960  # 固定渲染宽度
EMU_PER_INCH = 914400


def extract_render_paragraphs(slide):
    """
    按渲染顺序提取所有文本框段落的纯文本。
    返回 list[str]，与 HTML 中按 (shape_id, para_idx) 分组的 <span> 段落组严格 1:1 对齐。
    仅包含至少有一个 run 的段落（与渲染逻辑一致：无 run 的段落不产生 span）。
    """
    result = []

    def walk(shapes):
        for shape in shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(shape.shapes)
                elif shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        runs = para.runs
                        if not runs:
                            continue  # 无 run → 渲染时不产生 span，跳过保持对齐
                        txt = "".join(r.text or "" for r in runs)
                        result.append(txt)
            except Exception:
                continue

    walk(slide.shapes)
    return result


def _emu_to_px(emu_value, scale):
    """EMU → 像素（基于 scale）。"""
    if emu_value is None:
        return 0
    return round(emu_value * scale)


def _rgb_to_hex(rgb):
    """python-pptx RGB → hex 字符串。"""
    if rgb is None:
        return None
    try:
        return str(rgb)
    except Exception:
        return None


def _escape(text):
    """HTML 转义。"""
    if not text:
        return ""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")


def _text_frame_to_html(tf, shape_id, scale, left_px, top_px, width_px, height_px, transparent=False):
    """
    将 TextFrame 转为 HTML。
    返回带 data 属性的 <div> → <p> → <span> 嵌套结构。
    transparent=True 时文字透明（仅用于在 PNG 背景上做选中/高亮的隐形文字层）。
    """
    paragraphs_html = []
    for p_idx, para in enumerate(tf.paragraphs):
        runs_html = []
        for r_idx, run in enumerate(para.runs):
            font = run.font
            css_parts = []

            if font.size:
                # 字号: EMU → pt → CSS
                try:
                    pt_size = font.size / 12700  # EMU to points
                    css_parts.append(f"font-size:{pt_size:.1f}pt")
                except Exception:
                    pass
            if font.bold:
                css_parts.append("font-weight:bold")
            if font.italic:
                css_parts.append("font-style:italic")
            if transparent:
                # 透明文字层：文字不可见（PNG 背景已渲染），但可被选中/高亮
                css_parts.append("color:transparent")
            elif font.color and font.color.type is not None:
                try:
                    hex_color = _rgb_to_hex(font.color.rgb)
                    if hex_color and hex_color != "000000":
                        css_parts.append(f"color:#{hex_color}")
                except Exception:
                    pass  # 主题颜色等无法提取 RGB，跳过
            if font.name:
                css_parts.append(f'font-family:"{font.name}",sans-serif')

            style_str = ";".join(css_parts)
            attrs = [
                f'data-shape-id="{shape_id}"',
                f'data-para-idx="{p_idx}"',
                f'data-run-idx="{r_idx}"',
            ]
            if style_str:
                attrs.append(f'style="{style_str}"')

            text = _escape(run.text) if run.text else ""
            # 空文本需要占位以避免高亮定位失败
            if not text.strip():
                text = "&#8203;"  # zero-width space
            runs_html.append(f'<span {" ".join(attrs)}>{text}</span>')

        # 段落对齐
        align_css = ""
        if para.alignment == PP_ALIGN.CENTER:
            align_css = "text-align:center;"
        elif para.alignment == PP_ALIGN.RIGHT:
            align_css = "text-align:right;"
        elif para.alignment == PP_ALIGN.JUSTIFY:
            align_css = "text-align:justify;"

        p_style = f'style="{align_css}"' if align_css else ""
        paragraphs_html.append(f'<p {p_style}>{"".join(runs_html)}</p>')

    # 垂直对齐
    va_css = ""
    if tf.vertical_anchor is not None:
        va_map = {1: "top", 2: "middle", 3: "bottom"}
        va = va_map.get(tf.vertical_anchor, "top")
        va_css = f"display:flex;flex-direction:column;justify-content:{va};"

    word_wrap = "word-wrap:break-word;overflow-wrap:break-word;"
    return (
        f'<div style="{va_css}{word_wrap}width:100%;height:100%;">'
        f'{"".join(paragraphs_html)}'
        f'</div>'
    )


def _table_to_html(table):
    """将 PPT 表格转为 HTML <table>。"""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            ct = _escape(cell.text) if cell.text else "&nbsp;"
            cells.append(f'<td style="border:1px solid #ccc;padding:4px 8px;font-size:12px;">{ct}</td>')
        rows.append(f'<tr>{"".join(cells)}</tr>')
    return f'<table class="ppt-table" style="border-collapse:collapse;width:100%;height:100%;">{ "".join(rows)}</table>'


def _find_image_file(ppt_id, slide_idx, shape_id):
    """在图片目录中查找匹配的图片文件。"""
    from paths import data_dir
    images_dir = os.path.join(data_dir(), "uploads", ppt_id, "images")
    if not os.path.isdir(images_dir):
        return None

    # 精确匹配: slide{idx}_shape{id}.{ext}
    prefix = f"slide{slide_idx}_shape{shape_id}."
    for fname in os.listdir(images_dir):
        if fname.startswith(prefix):
            return f"/api/ppt/{ppt_id}/images/{fname}"
    return None


def _shape_to_html(shape, scale, ppt_id, slide_idx=0, transparent=False):
    """
    将单个 shape 转为绝对定位的 HTML div。
    transparent=True（高保真透明文字层）：图片/表格/图表不渲染（PNG 背景已含），
    仅渲染透明文字 span 供选中/高亮。
    """
    left_px = _emu_to_px(shape.left, scale)
    top_px = _emu_to_px(shape.top, scale)
    width_px = _emu_to_px(shape.width, scale)
    height_px = _emu_to_px(shape.height, scale)

    rotation_css = ""
    try:
        if shape.rotation:
            rotation_css = f"transform:rotate({shape.rotation}deg);"
    except Exception:
        pass

    shape_id = shape.shape_id
    inner = ""

    try:
        if shape.has_text_frame:
            try:
                inner = _text_frame_to_html(
                    shape.text_frame, shape_id, scale,
                    left_px, top_px, width_px, height_px, transparent=transparent
                )
            except Exception as e:
                # fallback: 直接提取纯文本
                color_css = "color:transparent;" if transparent else ""
                inner = f'<div style="width:100%;height:100%;display:flex;align-items:center;padding:8px;{color_css}">'
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    txt = _escape(para.text)
                    if txt.strip():
                        inner += f'<span data-shape-id="{shape_id}" data-para-idx="{pi}" data-run-idx="0">{txt}</span><br>'
                inner += '</div>'
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if transparent:
                # 高保真模式：图片已在 PNG 背景中，仅保留点击识图的透明热区
                inner = (
                    f'<img src="" data-shape-id="{shape_id}" '
                    f'data-image-path="shape{shape_id}" '
                    f'style="width:100%;height:100%;opacity:0;cursor:pointer;" alt="">'
                )
            else:
                img_src = _find_image_file(ppt_id, slide_idx, shape_id) or ""
                inner = (
                    f'<img src="{img_src}" data-shape-id="{shape_id}" '
                    f'data-image-path="shape{shape_id}" '
                    f'style="width:100%;height:100%;object-fit:contain;cursor:pointer;" '
                    f'alt="点击识图">'
                )
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            if shape.has_table and not transparent:
                inner = _table_to_html(shape.table)
            # transparent 模式下表格已在 PNG 中，不重复渲染
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # 递归渲染组合内的子 shape
            children = []
            try:
                for child in shape.shapes:
                    children.append(_shape_to_html(child, scale, ppt_id, slide_idx, transparent))
            except Exception:
                pass
            inner = "".join(children)
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            if not transparent:
                inner = f'<div style="display:flex;align-items:center;justify-content:center;height:100%;color:#999;font-size:12px;">📊 [图表]</div>'
            # transparent 模式下图表已在 PNG 中
        else:
            # 其他类型（如媒体、SmartArt 等）— 显示占位
            pass
    except Exception as e:
        inner = f'<!-- shape {shape_id} render error: {_escape(str(e))} -->'

    # 形状的背景/边框
    shape_style = f"position:absolute;left:{left_px}px;top:{top_px}px;width:{width_px}px;height:{height_px}px;overflow:visible;{rotation_css}"
    return f'<div style="{shape_style}" data-shape-id="{shape_id}">{inner}</div>'


def _render_fidelity_slide(ppt_id: str, slide_idx: int, span_data: dict, slide=None, emu_scale=None) -> str:
    """
    高保真渲染：PNG 背景 + 精确透明文字层 + 图片点击热区。

    文字 span 的 bbox 来自 PowerPoint 导出的 PDF（PyMuPDF 提取），与 PNG 同坐标系，
    因此选中/高亮区域与可见文字像素级对齐。

    每个 span 携带 data-shape-id / data-para-idx / data-run-idx（映射 PDF 的
    block/line/span 索引），供现有高亮系统定位，无需改动 highlight.js。
    图片用 python-pptx 的 EMU 坐标（显式存储，准确）叠加透明点击热区。
    """
    page_w = span_data.get("page_width", 960) or 960
    page_h = span_data.get("page_height", 540) or 540
    spans = span_data.get("spans", [])

    # PDF point → 显示像素：显示宽固定 DISPLAY_WIDTH_PX
    s = DISPLAY_WIDTH_PX / page_w
    disp_w = DISPLAY_WIDTH_PX
    disp_h = round(page_h * s)

    bg_css = (
        f"background-image:url('/api/ppt/{ppt_id}/render/slide{slide_idx}.png');"
        f"background-size:100% 100%;background-repeat:no-repeat;background-color:#fff;"
    )

    spans_html = []
    for sp in spans:
        left = sp["x"] * s
        top = sp["y"] * s
        width = sp["w"] * s
        height = sp["h"] * s
        # 字号按比例缩放，line-height 设为 bbox 高度让文字垂直填满，选区贴合
        font_px = sp["size"] * s
        text = _escape(sp["text"])
        spans_html.append(
            f'<span class="fz-span" '
            f'data-shape-id="{sp["block"]}" data-para-idx="{sp["line"]}" data-run-idx="{sp["span"]}" '
            f'style="position:absolute;left:{left:.2f}px;top:{top:.2f}px;'
            f'width:{width:.2f}px;height:{height:.2f}px;'
            f'font-size:{font_px:.2f}px;line-height:{height:.2f}px;'
            f'color:transparent;white-space:pre;overflow:hidden;">{text}</span>'
        )

    # 图片点击热区（识图用）— 用 python-pptx EMU 坐标
    hotspots_html = []
    if slide is not None and emu_scale is not None:
        def _collect_pics(shapes):
            for shape in shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        _collect_pics(shape.shapes)
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        l = _emu_to_px(shape.left, emu_scale)
                        t = _emu_to_px(shape.top, emu_scale)
                        w = _emu_to_px(shape.width, emu_scale)
                        h = _emu_to_px(shape.height, emu_scale)
                        hotspots_html.append(
                            f'<img src="" data-shape-id="{shape.shape_id}" '
                            f'data-image-path="shape{shape.shape_id}" '
                            f'style="position:absolute;left:{l}px;top:{t}px;width:{w}px;height:{h}px;'
                            f'opacity:0;cursor:pointer;" alt="点击识图">'
                        )
                except Exception:
                    continue
        try:
            _collect_pics(slide.shapes)
        except Exception:
            pass

    wrapper_style = (
        f"width:{disp_w}px;height:{disp_h}px;position:relative;{bg_css}"
    )
    return (
        f'<div class="slide-wrapper fidelity" style="{wrapper_style}" '
        f'data-slide-idx="{slide_idx}">{"".join(spans_html)}{"".join(hotspots_html)}</div>'
    )


def render_translated_fidelity_slide(ppt_id: str, slide_idx: int, translated_lines: list[str]) -> str:
    """
    渲染翻译后的高保真幻灯片：PNG 背景 + 可见的译文文字层。

    文字来自 spans.json（PowerPoint 导出 PDF → PyMuPDF 提取，100% 保真），
    按 (block, line) 分组后替换为译文。每个译文行用半透明白色背景覆盖
    PNG 中的原文，确保译文清晰可读。

    参数:
        ppt_id: PPT 唯一标识
        slide_idx: 0-based 幻灯片索引
        translated_lines: 译文行列表，与 fidelity spans 的 (block,line) 组一一对应
    """
    from render_service import get_slide_spans

    span_data = get_slide_spans(ppt_id, slide_idx)
    page_w = span_data.get("page_width", 960) or 960
    page_h = span_data.get("page_height", 540) or 540
    spans = span_data.get("spans", [])

    s = DISPLAY_WIDTH_PX / page_w
    disp_w = DISPLAY_WIDTH_PX
    disp_h = round(page_h * s)

    bg_css = (
        f"background-image:url('/api/ppt/{ppt_id}/render/slide{slide_idx}.png');"
        f"background-size:100% 100%;background-repeat:no-repeat;background-color:#fff;"
    )

    # 按 (block, line) 分组，与 _extract_fidelity_slides_lines 保持一致
    groups: dict[tuple, list[dict]] = {}
    for sp in spans:
        key = (sp["block"], sp["line"])
        groups.setdefault(key, []).append(sp)

    sorted_keys = sorted(groups.keys())
    spans_html = []

    for ti, key in enumerate(sorted_keys):
        line_spans = sorted(groups[key], key=lambda s: s["span"])

        # 源文本（用于判断该行是否确实有文字内容）
        source_text = "".join(s["text"] for s in line_spans).strip()
        if not source_text:
            continue  # 空行，无需覆盖

        # 整行的包围盒
        min_x = min(sp["x"] for sp in line_spans)
        min_y = min(sp["y"] for sp in line_spans)
        max_x = max(sp["x"] + sp["w"] for sp in line_spans)
        max_y = max(sp["y"] + sp["h"] for sp in line_spans)

        left = min_x * s
        top = min_y * s
        width = (max_x - min_x) * s
        height = (max_y - min_y) * s
        font_px = line_spans[0]["size"] * s

        zh = translated_lines[ti] if ti < len(translated_lines) else ""
        if zh:
            # 有译文：显示译文（半透明白底 + 黑色文字）
            zh_escaped = _escape(zh)
            spans_html.append(
                f'<span class="fz-span translated" '
                f'data-shape-id="{key[0]}" data-para-idx="{key[1]}" data-run-idx="0" '
                f'style="position:absolute;left:{left:.2f}px;top:{top:.2f}px;'
                f'min-width:{width:.2f}px;min-height:{height:.2f}px;'
                f'font-size:{font_px:.2f}px;line-height:{height:.2f}px;'
                f'color:#000;background:rgba(255,255,255,0.88);'
                f'padding:1px 3px;border-radius:2px;'
                f'white-space:nowrap;">{zh_escaped}</span>'
            )
        else:
            # 无译文：仅用白色遮罩覆盖 PNG 原文，避免英文残留
            spans_html.append(
                f'<span class="fz-span translated cover-only" '
                f'data-shape-id="{key[0]}" data-para-idx="{key[1]}" data-run-idx="0" '
                f'style="position:absolute;left:{left:.2f}px;top:{top:.2f}px;'
                f'width:{width:.2f}px;height:{height:.2f}px;'
                f'background:rgba(255,255,255,0.92);border-radius:2px;"></span>'
            )

    wrapper_style = f"width:{disp_w}px;height:{disp_h}px;position:relative;{bg_css}"
    return (
        f'<div class="slide-wrapper fidelity translated" style="{wrapper_style}" '
        f'data-slide-idx="{slide_idx}">{"".join(spans_html)}</div>'
    )


def render_slide(pres, slide_idx: int, ppt_id: str, fidelity: bool = False) -> str:
    """
    渲染单张幻灯片为 HTML 字符串。

    参数:
        pres: python-pptx Presentation 对象
        slide_idx: 0-based 幻灯片索引
        ppt_id: PPT 的唯一标识（用于图片 URL）
        fidelity: True=高保真模式（PNG 背景 + PDF 精确透明文字层），
                  False=普通模式（解析后的可见文字）

    返回:
        完整的幻灯片 HTML（含 slide-wrapper div）
    """
    if slide_idx < 0 or slide_idx >= len(pres.slides):
        return '<div class="slide-wrapper"><p>幻灯片索引超出范围</p></div>'

    slide = pres.slides[slide_idx]

    # 计算缩放比例
    slide_width_emu = pres.slide_width or 12192000  # 默认 10"
    slide_height_emu = pres.slide_height or 6858000  # 默认 7.5"
    scale = DISPLAY_WIDTH_PX / slide_width_emu
    slide_height_px = round(slide_height_emu * scale)

    # 高保真模式：PowerPoint 导出的 PNG 作背景 + PDF 提取的精确透明文字层
    if fidelity:
        from render_service import slide_image_path, get_slide_spans
        if slide_image_path(ppt_id, slide_idx):
            return _render_fidelity_slide(
                ppt_id, slide_idx, get_slide_spans(ppt_id, slide_idx), slide, scale
            )
        # PNG 不存在 → 降级到普通模式

    # 普通模式：幻灯片背景
    bg_css = "background-color:#ffffff;"
    try:
        bg = slide.background
        if bg.fill and bg.fill.type is not None:
            from pptx.enum.dml import MSO_FILL_TYPE
            try:
                if bg.fill.type == MSO_FILL_TYPE.SOLID and bg.fill.fore_color.rgb:
                    bg_css = f"background-color:#{_rgb_to_hex(bg.fill.fore_color.rgb)};"
            except Exception:
                pass
    except Exception:
        pass

    # 渲染各 shape
    shapes_html = []
    for shape in slide.shapes:
        try:
            shapes_html.append(_shape_to_html(shape, scale, ppt_id, slide_idx))
        except Exception as e:
            shapes_html.append(f'<!-- shape error: {_escape(str(e))} -->')

    wrapper_style = (
        f"width:{DISPLAY_WIDTH_PX}px;"
        f"height:{slide_height_px}px;"
        f"position:relative;"
        f"{bg_css}"
    )

    html = (
        f'<div class="slide-wrapper" style="{wrapper_style}" '
        f'data-slide-idx="{slide_idx}">'
        f'{"".join(shapes_html)}'
        f'</div>'
    )
    return html
